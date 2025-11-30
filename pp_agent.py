"""
create_presentation.py
----------------------
Generates a PowerPoint deck automatically using OpenAI GPT models and your company template.

Usage:
    python create_presentation.py --topic "DataCamp Overview" --slides 8 --template "CompanyTemplate.pptx"
"""

import argparse
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
import os
import logging
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
import json
from langchain_tavily import TavilySearch
from langchain.agents import create_agent

from dotenv import load_dotenv
load_dotenv()

# Configure logging
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('pp_agent.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# ---------- CONFIG ----------
MODEL = "gpt-4o-mini"  # fast + cost-effective
OUTPUT_FILE = "Generated_Presentation.pptx"
# ----------------------------



def analyze_slide_layouts(template_path: str) -> dict:
    """Analyze the template to understand available slide layouts and their structures."""
    logger.info(f"Analyzing slide layouts from template: {template_path}")
    
    try:
        prs = Presentation(template_path)
        logger.info(f"Template loaded successfully")
    except Exception as e:
        logger.error(f"Failed to load template: {str(e)}")
        raise
    
    layouts_info = []
    
    for idx, layout in enumerate(prs.slide_layouts):
        layout_data = {
            "index": idx,
            "name": layout.name,
            "placeholders": []
        }
        logger.debug(f"Layout {idx}: {layout.name}")
        
        # Analyze placeholders in this layout
        for shape in layout.placeholders:
            placeholder_info = {
                "idx": shape.placeholder_format.idx,
                "type": str(shape.placeholder_format.type),
                "name": shape.name if hasattr(shape, 'name') else "Unnamed"
            }
            layout_data["placeholders"].append(placeholder_info)
            logger.debug(f"  - Placeholder: {placeholder_info['name']} (type: {placeholder_info['type']})")
        
        layouts_info.append(layout_data)
    
    logger.info(f"Found {len(layouts_info)} layouts in template")
    return {
        "layouts": layouts_info,
        "total_layouts": len(layouts_info)
    }

def generate_slide_outline(topic: str, n_slides: int, instructions: str, layouts_info: dict):
    """Generate a slide outline using LangChain + OpenAI with web search tool."""
    logger.info(f"Generating slide outline for topic: {topic}")
    logger.info(f"Requested slides: {n_slides}")
    logger.info(f"Available layouts: {layouts_info['total_layouts']}")
    
    # Initialize search tool
    try:
        search_tool = TavilySearch(max_results=5,topic="general")
        logger.info("Tavily search tool initialized successfully")
    except Exception as e:
        logger.warning(f"Failed to initialize Tavily search: {str(e)}, continuing without search")
        search_tool = None

    # Initialize LLM with tool binding
    try:
        llm = ChatOpenAI(
            model="gpt-4o",
            temperature=1,
           # model_kwargs={"response_format": {"type": "json_object"}}
        )
        
        # Bind the search tool to the LLM if available
        if search_tool:
            agent = create_agent(llm, [search_tool])
            logger.info("LLM initialized with search tool binding")
        else:
            logger.error("Failed to initialize LLM: No Seatch tool available")
    except Exception as e:
        logger.error(f"Failed to initialize LLM: {str(e)}")
        raise

    # Format layouts information for the prompt
    layouts_description = "\n".join([
        f"- Layout {l['index']}: '{l['name']}' with placeholders: {', '.join([p['name'] for p in l['placeholders']])}"
        for l in layouts_info['layouts']
    ])
    
    prompt_template = ChatPromptTemplate.from_template("""
    You are creating a professional internal PowerPoint presentation about "{topic}".
    
    IMPORTANT: If you need current information, statistics, or facts about this topic, use the tavily_search_results_json tool to search the web first. This will help make the presentation accurate and current.
    
    After gathering information (if needed), produce {n_slides} slides in **JSON** format.

    The PowerPoint template has the following slide layouts available:
    {layouts_description}

    The user has provided the following detailed instructions for what they want in the slides:
    ---
    {instructions}
    ---
    
    Ensure every slide aligns with these instructions and uses accurate, current information.

    For each slide, you must:
    1. Select the most appropriate layout_index based on the content type
    2. Provide content that matches the placeholders in that layout
    3. Include a title
    4. Include bullet points or content as appropriate for the layout
    5. Add speaker notes

    Each slide must include:
    - "layout_index": the index number of the layout to use (choose wisely based on content)
    - "title": short slide title
    - "content": either a list of bullet points OR a single text block, depending on the layout
    - "notes": 2–3 sentence speaker notes explaining how to present the slide

    Output format example:
    {{
      "slides": [
        {{
          "layout_index": 2,
          "title": "What is DataCamp?",
          "content": ["Online platform for data skills", "Python, R, SQL, Power BI", "Used by 10M+ learners"],
          "notes": "Introduce DataCamp as a flexible platform..."
        }},
        {{
          "layout_index": 3,
          "title": "Our Approach",
          "content": ["Hands-on learning", "Real-world projects", "Expert instructors"],
          "notes": "Explain the unique teaching methodology..."
        }}
      ]
    }}

    Respond ONLY with valid JSON wrapped in a "slides" array.
    """)

    chain = prompt_template | agent

    logger.info("Invoking LLM to generate slide content (with web search capability)...")
    try:
        result = chain.invoke({
            "topic": topic,
            "n_slides": n_slides,
            "instructions": instructions,
            "layouts_description": layouts_description
        })
        logger.info("LLM invocation completed")
    except Exception as e:
        logger.error(f"LLM invocation failed: {str(e)}")
        raise

    #content = result.content
    content = result['messages'][-1].content # Get the last message content
    logger.debug(f"Raw LLM response length: {len(content)} characters")

    # Parse JSON output safely
    try:
        parsed_data = json.loads(content.replace('```json','').replace('```',''))
        logger.info(f"Successfully parsed JSON with {len(parsed_data.get('slides', []))} slides")
        return parsed_data
    except Exception as e:
        logger.error(f"Error parsing JSON: {str(e)}")
        logger.error(f"Raw content: {content[:500]}...")  # Log first 500 chars
        return {"slides": []}


def build_presentation(slides_data: dict, template_path: str, output_path: str):
    """Populate slides into a PowerPoint using the company template with dynamic layout selection."""
    logger.info(f"Building presentation from template: {template_path}")
    logger.info(f"Output will be saved to: {output_path}")
    
    try:
        prs = Presentation(template_path)
        logger.info("Template loaded for building")
    except Exception as e:
        logger.error(f"Failed to load template for building: {str(e)}")
        raise
    
    slides = slides_data.get('slides', [])
    logger.info(f"Processing {len(slides)} slides")
    
    for idx, slide_data in enumerate(slides, 1):
        logger.info(f"Processing slide {idx}/{len(slides)}: {slide_data.get('title', 'Untitled')}")
        
        layout_index = slide_data.get('layout_index', 1)
        logger.debug(f"  Using layout index: {layout_index}")
        
        # Ensure layout index is valid
        if layout_index >= len(prs.slide_layouts):
            logger.warning(f"Layout index {layout_index} not found, using layout 1")
            layout_index = 1
        
        # Add slide with selected layout
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
            logger.debug(f"  Slide added successfully")
        except Exception as e:
            logger.error(f"Failed to add slide with layout {layout_index}: {str(e)}")
            continue
        
        # Get text shapes
        text_shapes = [s for s in slide.shapes if s.has_text_frame]
        logger.debug(f"  Found {len(text_shapes)} text shapes")
        
        if not text_shapes:
            logger.warning(f"No text shapes found in layout {layout_index}")
            continue
        
        # Set title (usually first text shape)
        if len(text_shapes) > 0:
            title_text = slide_data.get('title', 'Untitled')
            text_shapes[0].text = title_text
            logger.debug(f"  Title set: {title_text}")
        
        # Set content (usually second text shape if it exists)
        if len(text_shapes) > 1:
            content = slide_data.get('content', [])
            logger.debug(f"  Setting content (type: {type(content).__name__})")
            
            try:
                body = text_shapes[1].text_frame
                body.clear()
                
                # Handle both list and string content
                if isinstance(content, list):
                    logger.debug(f"  Adding {len(content)} bullet points")
                    for item in content:
                        p = body.add_paragraph()
                        p.text = str(item)
                        p.level = 0
                        if p.font:
                            p.font.size = Pt(18)
                else:
                    logger.debug(f"  Adding text content")
                    body.text = str(content)
            except Exception as e:
                logger.error(f"Failed to set content: {str(e)}")
        
        # Add speaker notes
        notes = slide_data.get('notes', '')
        if notes:
            try:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_frame.text = notes
                logger.debug(f"  Speaker notes added")
            except Exception as e:
                logger.error(f"Failed to add speaker notes: {str(e)}")

    logger.info(f"Saving presentation to: {output_path}")
    try:
        prs.save(output_path)
        logger.info(" Presentation saved successfully")
        print(f" Presentation saved to: {output_path}")
    except Exception as e:
        logger.error(f"Failed to save presentation: {str(e)}")
        raise

def main():
    logger.info("=== PowerPoint Builder Started ===")
    
    parser = argparse.ArgumentParser(description="Generate PowerPoint slides with OpenAI.")
    parser.add_argument("--topic", required=True, help="Presentation topic, e.g., 'DataCamp Overview'")
    parser.add_argument("--instructions", type=str, default="Make it professional and suitable for an internal company presentation.", help="Additional instructions for slide content")
    parser.add_argument("--slides", type=int, default=8, help="Number of slides to generate")
    parser.add_argument("--template", required=True, help="Path to company PowerPoint template")
    args = parser.parse_args()
    
    logger.info(f"Arguments received:")
    logger.info(f"  Topic: {args.topic}")
    logger.info(f"  Slides: {args.slides}")
    logger.info(f"  Template: {args.template}")
    logger.info(f"  Instructions: {args.instructions[:100]}...")

    try:
        # First, analyze the template to understand available layouts
        print(" Analyzing template layouts...")
        layouts_info = analyze_slide_layouts(args.template)
        print(f"Found {layouts_info['total_layouts']} slide layouts")
        
        # Generate slides with layout awareness
        print(f" Generating {args.slides} slides about '{args.topic}'...")
        slides_data = generate_slide_outline(args.topic, args.slides, args.instructions, layouts_info)
        
        if not slides_data.get('slides'):
            logger.error("No slides generated by LLM")
            print(" Failed to generate slides")
            return 1
        
        # Build the presentation
        print(" Building presentation...")
        build_presentation(slides_data, args.template, OUTPUT_FILE)
        
        logger.info("=== PowerPoint Builder Completed Successfully ===")
        return 0
        
    except Exception as e:
        logger.exception(f"Fatal error in main: {str(e)}")
        print(f" Error: {str(e)}")
        return 1


if __name__ == "__main__":
    exit(main())